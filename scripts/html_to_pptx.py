import os
import re
import sys
import json
import argparse
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# Mapping text alignments
ALIGN_MAP = {
    'left': PP_ALIGN.LEFT,
    'center': PP_ALIGN.CENTER,
    'right': PP_ALIGN.RIGHT,
    'justify': PP_ALIGN.JUSTIFY
}

# Mapping chart types to python-pptx chart types
CHART_TYPE_MAP = {
    'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
    'bar': XL_CHART_TYPE.BAR_CLUSTERED,
    'line': XL_CHART_TYPE.LINE,
    'pie': XL_CHART_TYPE.PIE,
    'area': XL_CHART_TYPE.AREA,
}

def parse_dimension(style_str, prop_name):
    """
    Parses a dimension from CSS inline style.
    Example style_str: 'top: 0.8in; left: 1.0in; width: 11.33in; height: 1.2in;'
    Returns the python-pptx measurement (Inches, Cm, Pt).
    """
    pattern = r'(?:^|;\s*)' + re.escape(prop_name) + r'\s*:\s*([\d.]+)\s*(in|cm|pt|px)?'
    match = re.search(pattern, style_str, re.IGNORECASE)
    if not match:
        return None
    val = float(match.group(1))
    unit = match.group(2)
    if unit:
        unit = unit.lower()
    else:
        unit = 'in' # Default to inches
        
    if unit == 'in':
        return Inches(val)
    elif unit == 'cm':
        return Cm(val)
    elif unit == 'pt':
        return Pt(val)
    elif unit == 'px':
        # Assume 96 DPI: 1 inch = 96 pixels
        return Inches(val / 96.0)
    return Inches(val)

def parse_color(color_str):
    """
    Parses color string like '#2F5597' or 'rgb(47, 85, 151)' to RGBColor.
    """
    if not color_str:
        return None
    color_str = color_str.strip()
    if color_str.startswith('#'):
        hex_val = color_str.lstrip('#')
        if len(hex_val) == 3:
            hex_val = ''.join([c*2 for c in hex_val])
        return RGBColor(int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16))
    elif color_str.startswith('rgb'):
        nums = re.findall(r'\d+', color_str)
        if len(nums) >= 3:
            return RGBColor(int(nums[0]), int(nums[1]), int(nums[2]))
    return None

def apply_style_to_run(run, style_dict):
    """
    Helper to apply font properties to a text run.
    """
    if style_dict.get('size'):
        run.font.size = style_dict['size']
    if style_dict.get('name'):
        run.font.name = style_dict['name']
    if style_dict.get('color'):
        run.font.color.rgb = style_dict['color']
    if style_dict.get('bold') is not None:
        run.font.bold = style_dict['bold']
    if style_dict.get('italic') is not None:
        run.font.italic = style_dict['italic']

def append_runs_from_html(paragraph, element, style):
    """
    Recursively processes HTML inline formatting (strong, em, code, etc.) 
    and adds formatted runs to a paragraph.
    """
    for child in element.children:
        if isinstance(child, str):
            text = child
            # Only add if there is non-whitespace text or single spaces
            if text.strip() or text == ' ':
                run = paragraph.add_run()
                run.text = text
                apply_style_to_run(run, style)
        elif child.name in ['strong', 'b']:
            new_style = style.copy()
            new_style['bold'] = True
            append_runs_from_html(paragraph, child, new_style)
        elif child.name in ['em', 'i']:
            new_style = style.copy()
            new_style['italic'] = True
            append_runs_from_html(paragraph, child, new_style)
        elif child.name == 'span':
            new_style = style.copy()
            # Try to parse inline color from span style
            span_style = child.get('style', '')
            color_match = re.search(r'color\s*:\s*([^;]+)', span_style)
            if color_match:
                parsed_c = parse_color(color_match.group(1))
                if parsed_c:
                    new_style['color'] = parsed_c
            append_runs_from_html(paragraph, child, new_style)
        elif child.name == 'code':
            new_style = style.copy()
            new_style['name'] = 'Courier New'
            append_runs_from_html(paragraph, child, new_style)
        elif child.name is not None:
            # For general elements, just recurse
            append_runs_from_html(paragraph, child, style)

def process_text_box(slide, el):
    """
    Parses and adds an editable Text Frame from .pptx-text element.
    """
    style_str = el.get('style', '')
    left = parse_dimension(style_str, 'left')
    top = parse_dimension(style_str, 'top')
    width = parse_dimension(style_str, 'width')
    height = parse_dimension(style_str, 'height')
    
    if left is None or top is None:
        print(f"Warning: Text box missing left or top coordinates: {el}")
        return
        
    width = width if width is not None else Inches(5.0)
    height = height if height is not None else Inches(1.0)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Read data attributes for baseline font styles
    default_font = {
        'size': Pt(int(el.get('data-font-size'))) if el.get('data-font-size') else None,
        'name': el.get('data-font-family') or el.get('data-font-name'),
        'color': parse_color(el.get('data-font-color')) if el.get('data-font-color') else None,
        'bold': el.get('data-bold', '').lower() == 'true',
        'italic': el.get('data-italic', '').lower() == 'true'
    }
    
    align_str = el.get('data-align')
    align = ALIGN_MAP.get(align_str.lower()) if align_str else None
    
    # Check for bullet points (ul/ol) or sub-paragraphs inside the element
    blocks = el.find_all(['ul', 'ol', 'p', 'div'], recursive=False)
    
    if not blocks:
        p = tf.paragraphs[0]
        if align:
            p.alignment = align
        append_runs_from_html(p, el, default_font)
    else:
        is_first = True
        for block in el.children:
            if isinstance(block, str):
                if not block.strip():
                    continue
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                if align:
                    p.alignment = align
                run = p.add_run()
                run.text = block
                apply_style_to_run(run, default_font)
            elif block.name in ['ul', 'ol']:
                # List items
                for li in block.find_all('li'):
                    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                    is_first = False
                    p.level = 0
                    
                    # Prepend custom bullet for non-placeholder text frames
                    bullet_run = p.add_run()
                    bullet_run.text = '• '
                    apply_style_to_run(bullet_run, default_font)
                    
                    if align:
                        p.alignment = align
                    append_runs_from_html(p, li, default_font)
            elif block.name in ['p', 'div']:
                p = tf.paragraphs[0] if is_first else tf.add_paragraph()
                is_first = False
                if align:
                    p.alignment = align
                append_runs_from_html(p, block, default_font)

def process_table(slide, el):
    """
    Parses and adds an editable PPTX Table from table.pptx-table element.
    """
    style_str = el.get('style', '')
    left = parse_dimension(style_str, 'left')
    top = parse_dimension(style_str, 'top')
    width = parse_dimension(style_str, 'width')
    height = parse_dimension(style_str, 'height')
    
    if left is None or top is None:
        print(f"Warning: Table missing left or top coordinates: {el}")
        return
        
    width = width if width is not None else Inches(8.0)
    height = height if height is not None else Inches(2.0)
    
    rows_els = el.find_all('tr')
    num_rows = len(rows_els)
    if num_rows == 0:
        return
        
    num_cols = max(len(row.find_all(['td', 'th'])) for row in rows_els)
    if num_cols == 0:
        return
        
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table
    
    # Set explicit column widths if provided
    col_widths_str = el.get('data-col-widths')
    if col_widths_str:
        try:
            widths = json.loads(col_widths_str)
            for idx, w in enumerate(widths):
                if idx < len(table.columns):
                    table.columns[idx].width = Inches(w)
        except Exception as e:
            print(f"Error parsing table col widths: {e}")
            
    default_font_size = Pt(int(el.get('data-font-size'))) if el.get('data-font-size') else Pt(14)
    
    for r_idx, row_el in enumerate(rows_els):
        row_style = row_el.get('style', '')
        bg_match = re.search(r'background-color\s*:\s*([^;]+)', row_style)
        row_bg_color = parse_color(bg_match.group(1)) if bg_match else None
        
        cells = row_el.find_all(['td', 'th'])
        for c_idx, cell_el in enumerate(cells):
            if c_idx >= num_cols:
                break
                
            cell = table.cell(r_idx, c_idx)
            cell.text = "" # Clear defaults
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # cell-level background color overrides row background
            cell_style = cell_el.get('style', '')
            cell_bg_match = re.search(r'background-color\s*:\s*([^;]+)', cell_style)
            cell_bg_color = parse_color(cell_bg_match.group(1)) if cell_bg_match else row_bg_color
            
            if cell_bg_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_bg_color
                
            # Parse text color inside cell styling or row styling
            color_match = re.search(r'color\s*:\s*([^;]+)', cell_style + ';' + row_style)
            cell_text_color = parse_color(color_match.group(1)) if color_match else None
            
            is_header = cell_el.name == 'th'
            cell_font = {
                'size': default_font_size,
                'name': 'Arial',
                'color': cell_text_color or (RGBColor(255, 255, 255) if is_header else RGBColor(0, 0, 0)),
                'bold': is_header,
                'italic': False
            }
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            
            append_runs_from_html(p, cell_el, cell_font)

def process_chart(slide, el):
    """
    Parses and adds an editable native Excel-backed Chart from .pptx-chart element.
    """
    style_str = el.get('style', '')
    left = parse_dimension(style_str, 'left')
    top = parse_dimension(style_str, 'top')
    width = parse_dimension(style_str, 'width')
    height = parse_dimension(style_str, 'height')
    
    if left is None or top is None:
        print(f"Warning: Chart missing left or top coordinates: {el}")
        return
        
    width = width if width is not None else Inches(5.0)
    height = height if height is not None else Inches(4.0)
    
    chart_type_str = el.get('data-chart-type', 'column').lower()
    chart_type = CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    categories_str = el.get('data-chart-categories')
    series_str = el.get('data-chart-series')
    title_str = el.get('data-chart-title')
    
    if not categories_str or not series_str:
        print(f"Warning: Chart missing categories or series data: {el}")
        return
        
    try:
        categories = json.loads(categories_str)
        series_data = json.loads(series_str)
    except Exception as e:
        print(f"Error parsing chart JSON data: {e}")
        return
        
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for s in series_data:
        name = s.get('name', 'Series')
        values = s.get('values', [])
        chart_data.add_series(name, tuple(values))
        
    chart_shape = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_shape.chart
    
    chart.has_legend = True
    if title_str:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_str

def html_to_pptx(html_path, pptx_path):
    """
    Reads an HTML file, parses the slide containers, and writes a PPTX presentation.
    """
    if not os.path.exists(html_path):
        print(f"Error: HTML file '{html_path}' not found.")
        sys.exit(1)
        
    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f.read(), 'lxml')
        
    prs = Presentation()
    # Set default widescreen 16:9 layout (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide
    
    slide_elements = soup.find_all(class_='slide')
    print(f"Found {len(slide_elements)} slides in HTML.")
    
    for idx, slide_el in enumerate(slide_elements):
        print(f"Processing slide {idx+1}/{len(slide_elements)}...")
        slide = prs.slides.add_slide(blank_layout)
        
        # 1. Process Text Boxes
        text_boxes = slide_el.find_all(class_='pptx-text')
        for tb in text_boxes:
            process_text_box(slide, tb)
            
        # 2. Process Tables (either table tag or class pptx-table)
        tables = slide_el.find_all(class_='pptx-table')
        # Also process generic table tags that might not have class pptx-table but are inside the slide
        table_tags = slide_el.find_all('table')
        all_tables = list(set(tables + table_tags))
        for tbl in all_tables:
            process_table(slide, tbl)
            
        # 3. Process Charts
        charts = slide_el.find_all(class_='pptx-chart')
        for crt in charts:
            process_chart(slide, crt)
            
    prs.save(pptx_path)
    print(f"Success: Presentation saved to '{pptx_path}'")

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description="Convert structured HTML to native editable PPTX.")
    parser.add_argument("html_input", help="Path to input HTML file")
    parser.add_argument("pptx_output", nargs="?", help="Path to output PPTX file (optional)")
    
    args = parser.parse_args()
    
    html_input = args.html_input
    pptx_output = args.pptx_output
    if not pptx_output:
        pptx_output = os.path.splitext(html_input)[0] + ".pptx"
        
    html_to_pptx(html_input, pptx_output)
